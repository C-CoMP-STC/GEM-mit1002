#!GEM-mit1002/.venv/bin/python
"""
This script creates a PowerPoint presentation with slides containing Escher map
images. It uses a predefined template and layout to match the C-CoMP
presentation template, and adds images along with titles and species names.
"""
import os

from pptx import Presentation
from pptx.util import Inches

# --- CONFIGURATION ---

# Define path names
FILE_DIR = os.path.dirname(os.path.abspath(__file__))
PNG_DIR = os.path.join(FILE_DIR, "results", "escher_plots", "png")

# Define image positions for reuse
IMG_POSITIONS = {
    "amac_left": {"left": Inches(-0.58), "top": Inches(1.63), "width": Inches(7.25)},
    "ecoli_right": {"left": Inches(6.67), "top": Inches(1.38), "width": Inches(6.67)},
    "amac_right": {"left": Inches(6.67), "top": Inches(1.63), "width": Inches(7.25)},
}

# Define the exact order of slides to be generated
SLIDE_ORDER = [
    "Glucose+Ammonia",
    "Glucose+Nitrate",
    "Glucose+Ammonia+Nitrate",
    "Acetate+Ammonia",
    "Acetate+Nitrate",
    "Acetate+Ammonia+Nitrate",
    "Glucose+Acetate+Ammonia+Nitrate",
]

# --- HELPER FUNCTIONS ---


def format_source_string(source_string: str) -> str:
    """
    Formats a '+' separated string into a human-readable list for titles.
    e.g., "A+B+C" becomes "A, B & C".
    """
    parts = source_string.split("+")
    if len(parts) <= 1:
        return source_string
    if len(parts) == 2:
        return f"{parts[0]} & {parts[1]}"
    # For 3 or more items, comma-separate all but the last one.
    return f"{', '.join(parts[:-1])} & {parts[-1]}"


def create_comparison_slides(
    prs,
    layout,
    base_prefix,
    partner_prefix,
    left_pos_key,
    right_pos_key,
    left_text,
    right_text,
):
    """
    Finds image pairs and generates a 'growth' and 'blockages' slide for each C/N source.

    Args:
        prs: The Presentation object.
        layout: The slide layout to use.
        base_prefix: Filename prefix for the "left" image.
        partner_prefix: Filename prefix for the "right" image.
        left_pos_key: Key for the left image position in IMG_POSITIONS.
        right_pos_key: Key for the right image position in IMG_POSITIONS.
        left_text: Text for the left species label.
        right_text: Text for the right species label.
    """
    print(f"\n--- Generating slides for {base_prefix} vs {partner_prefix} ---")
    for cn_source in SLIDE_ORDER:
        for slide_type in ["growth", "blockages"]:
            # Format the title string for readability
            formatted_title_source = format_source_string(cn_source)

            # Determine filenames based on slide type
            if slide_type == "growth":
                base_filename = f"{base_prefix}{cn_source}.png"
                partner_filename = f"{partner_prefix}{cn_source}.png"
                title = f"Growth on {formatted_title_source}"
            else:  # blockages
                base_filename = f"{base_prefix}{cn_source}_blocked_reactions.png"
                partner_filename = f"{partner_prefix}{cn_source}_blocked_reactions.png"
                title = f"Blockages on {formatted_title_source}"

            base_filepath = os.path.join(PNG_DIR, base_filename)
            partner_filepath = os.path.join(PNG_DIR, partner_filename)

            # Check if both image files exist before creating a slide
            if os.path.exists(base_filepath) and os.path.exists(partner_filepath):
                print(f"  > Creating '{slide_type}' slide for {cn_source}")
                slide = prs.slides.add_slide(layout)

                # Add pictures
                slide.shapes.add_picture(base_filepath, **IMG_POSITIONS[left_pos_key])
                slide.shapes.add_picture(
                    partner_filepath, **IMG_POSITIONS[right_pos_key]
                )

                # Set title and text boxes
                slide.shapes[0].text_frame.text = title
                slide.shapes[2].text_frame.text = left_text
                slide.shapes[3].text_frame.text = right_text
            else:
                print(f"  > SKIPPING '{slide_type}' for {cn_source} (missing images)")


# --- MAIN SCRIPT ---


def main():
    """
    Generates a PowerPoint presentation from Escher map images.
    """
    # Open template and choose layout
    prs = Presentation(os.path.join(FILE_DIR, "slide-template.pptx"))
    my_layout = prs.slide_layouts[9]

    # --- Generate Slides for Each Comparison ---

    # 1. AMAC ORIGINAL VS E COLI
    create_comparison_slides(
        prs,
        my_layout,
        base_prefix="amac_Original_",
        partner_prefix="ecoli_",
        left_pos_key="amac_left",
        right_pos_key="ecoli_right",
        left_text="A. macleodii MIT1002",
        right_text="E. coli (iJO1366)",
    )

    # 2. AMAC ORIGINAL VS LUMPED REACTIONS
    create_comparison_slides(
        prs,
        my_layout,
        base_prefix="amac_Original_",
        partner_prefix="amac_With_Lumped_Reactions_",
        left_pos_key="amac_left",
        right_pos_key="amac_right",
        left_text="MIT1002 (As is)",
        right_text="MIT1002 (With Lumped Reactions)",
    )

    # 3. AMAC STRICT ATP VS STRICT ATP + NUCLEOTIDE BALANCING
    create_comparison_slides(
        prs,
        my_layout,
        base_prefix="amac_Strict_ATP_Production_",
        partner_prefix="amac_Strict_ATP_Production_Nucleotide_Balancing_",
        left_pos_key="amac_left",
        right_pos_key="amac_right",
        left_text="MIT1002 (Strict ATP Production)",
        right_text="MIT1002 (Strict ATP + Nucleotide Balancing)",
    )

    # 4. AMAC STRICT ATP VS STRICT ATP + LUMPED REACTIONS
    create_comparison_slides(
        prs,
        my_layout,
        base_prefix="amac_Strict_ATP_Production_",
        partner_prefix="amac_Strict_ATP_Production_With_Lumped_Reactions_",
        left_pos_key="amac_left",
        right_pos_key="amac_right",
        left_text="MIT1002 (Strict ATP Production)",
        right_text="MIT1002 (Strict ATP + Lumped Reactions)",
    )

    # Save the output file
    output_path = os.path.join(FILE_DIR, "results", "generated-slides.pptx")
    prs.save(output_path)
    print(f"\nPresentation saved to {output_path}")


if __name__ == "__main__":
    main()
